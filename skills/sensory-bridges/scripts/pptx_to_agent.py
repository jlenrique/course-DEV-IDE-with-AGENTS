"""PPTX sensory bridge — deterministic slide text and structure extraction.

Primary path for G3 text verification: extracts exact text objects from
PPTX slide data, not OCR from rendered images. Every text string returned
is the literal value Gamma placed in the presentation.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches  # noqa: F401 — imported for future layout analysis

from skills.sensory_bridges.scripts.bridge_utils import build_response

logger = logging.getLogger(__name__)


def _extract_text_frames(slide) -> list[str]:
    """Extract all text from a slide's shapes."""
    texts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = "\n".join(
                paragraph.text for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            )
            if full_text.strip():
                texts.append(full_text.strip())
    return texts


def _extract_image_refs(slide) -> list[dict[str, str]]:
    """Extract image references from a slide."""
    refs: list[dict[str, str]] = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            image = shape.image
            refs.append({
                "content_type": image.content_type,
                "filename": image.filename if hasattr(image, "filename") else "",
                "width_px": shape.width,
                "height_px": shape.height,
            })
    return refs


def _extract_notes(slide) -> str:
    """Extract speaker notes from a slide."""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""


def extract_pptx(
    artifact_path: str | Path,
    gate: str = "G3",
    **kwargs: Any,
) -> dict[str, Any]:
    """Extract structured content from a PPTX file.

    Args:
        artifact_path: Path to the .pptx file.
        gate: Production gate identifier.

    Returns:
        Canonical perception response with slides[], total_slides, confidence.
    """
    path = Path(artifact_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {path}")

    try:
        prs = Presentation(str(path))
    except Exception as e:
        return build_response(
            modality="pptx",
            artifact_path=path,
            confidence="LOW",
            confidence_rationale=f"Failed to parse PPTX: {e}",
            slides=[],
            total_slides=0,
        )

    slides_data: list[dict[str, Any]] = []
    has_issues = False

    for i, slide in enumerate(prs.slides, start=1):
        try:
            text_frames = _extract_text_frames(slide)
            image_refs = _extract_image_refs(slide)
            notes = _extract_notes(slide)

            slides_data.append({
                "slide_number": i,
                "text_frames": text_frames,
                "image_refs": image_refs,
                "notes": notes,
            })
        except Exception as e:
            logger.warning("Error extracting slide %d: %s", i, e)
            slides_data.append({
                "slide_number": i,
                "text_frames": [],
                "image_refs": [],
                "notes": "",
                "extraction_error": str(e),
            })
            has_issues = True

    total = len(slides_data)

    if total == 0:
        confidence = "LOW"
        rationale = "PPTX contains no slides"
    elif has_issues:
        confidence = "MEDIUM"
        rationale = f"Extracted {total} slides but encountered errors on some"
    else:
        all_have_text = all(s["text_frames"] for s in slides_data)
        confidence = "HIGH" if all_have_text else "MEDIUM"
        rationale = (
            f"Extracted {total} slides, all with text content"
            if all_have_text
            else f"Extracted {total} slides, some without text frames"
        )

    return build_response(
        modality="pptx",
        artifact_path=path,
        confidence=confidence,
        confidence_rationale=rationale,
        slides=slides_data,
        total_slides=total,
    )
